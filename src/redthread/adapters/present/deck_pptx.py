"""Render a slide deck from a list of curated handoffs, via python-pptx."""

from pptx import Presentation
from pptx.util import Pt

from redthread.models import Handoff

_TITLE_LAYOUT = 0
_TITLE_AND_CONTENT_LAYOUT = 1


def render_deck(handoffs: list[Handoff], run_id: str) -> Presentation:
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[_TITLE_LAYOUT])
    title_slide.shapes.title.text = f"Run {run_id}"
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = f"{len(handoffs)} phase(s) summarized"

    for handoff in handoffs:
        slide = prs.slides.add_slide(prs.slide_layouts[_TITLE_AND_CONTENT_LAYOUT])
        slide.shapes.title.text = f"{handoff.from_phase}: {handoff.headline}"

        body = slide.placeholders[1].text_frame
        body.clear()
        first = True
        for key, value in handoff.key_results.items():
            para = body.paragraphs[0] if first else body.add_paragraph()
            para.text = f"{key}: {value}"
            para.font.size = Pt(18)
            first = False
        for decision in handoff.decisions:
            para = body.paragraphs[0] if first else body.add_paragraph()
            para.text = f"Decision: {decision}"
            para.font.size = Pt(16)
            first = False

    return prs
