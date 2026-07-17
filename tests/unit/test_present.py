from pptx import Presentation

from redthread.adapters.present.deck_pptx import render_deck
from redthread.adapters.present.docs_site import render_docs_tree
from redthread.adapters.present.report_md import render_report
from redthread.models import Handoff

_H1 = Handoff(
    from_phase="build",
    run_id="r1",
    headline="build ok, 0 warnings",
    key_results={"warnings": 0},
    decisions=["enabled tree-shaking"],
)
_H2 = Handoff(
    from_phase="test",
    run_id="r1",
    headline="tests passed, 87% coverage",
    key_results={"coverage_pct": 87},
    open_questions=["flaky test on CI?"],
)


def test_render_report_contains_headline_metrics_and_decisions():
    report = render_report([_H1, _H2], "r1")
    assert "build ok, 0 warnings" in report
    assert "tests passed, 87% coverage" in report
    assert "| warnings | 0 |" in report
    assert "enabled tree-shaking" in report
    assert "flaky test on CI?" in report


def test_render_report_handles_no_handoffs():
    report = render_report([], "r1")
    assert "r1" in report
    assert "No upstream" in report


def test_render_deck_has_one_title_slide_plus_one_per_handoff():
    prs = render_deck([_H1, _H2], "r1")
    assert len(prs.slides) == 3  # title + 2 handoffs

    slide_texts = [s.shapes.title.text for s in prs.slides]
    assert slide_texts[0] == "Run r1"
    assert "build: build ok, 0 warnings" in slide_texts[1]
    assert "test: tests passed, 87% coverage" in slide_texts[2]


def test_render_deck_body_contains_key_results_and_decisions(tmp_path):
    prs = render_deck([_H1], "r1")
    deck_path = tmp_path / "deck.pptx"
    prs.save(str(deck_path))

    reopened = Presentation(str(deck_path))
    body_slide = list(reopened.slides)[1]
    body_text = "\n".join(p.text for p in body_slide.placeholders[1].text_frame.paragraphs)
    assert "warnings: 0" in body_text
    assert "Decision: enabled tree-shaking" in body_text


def test_render_docs_tree_writes_index_and_one_file_per_phase(tmp_path):
    written = render_docs_tree([_H1, _H2], "r1", tmp_path / "docs")
    assert len(written) == 3  # index + build + test

    index = (tmp_path / "docs" / "index.md").read_text(encoding="utf-8")
    assert "build" in index and "test" in index

    build_page = (tmp_path / "docs" / "build.md").read_text(encoding="utf-8")
    assert "build ok, 0 warnings" in build_page
    assert "warnings" in build_page
