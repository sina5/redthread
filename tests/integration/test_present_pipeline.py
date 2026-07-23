"""The present phase, running on a THIRD node, renders a report + deck +
docs tree whose facts match the upstream handoffs exactly — for both
example pipelines, proving the present layer is domain-neutral too.
"""

import subprocess

from pptx import Presentation

from redthread.adapters.examples.app_build import run_build
from redthread.adapters.examples.app_test import run_tests
from redthread.adapters.examples.ml_eval import run_eval
from redthread.adapters.examples.ml_train import run_training
from redthread.adapters.present import run_present
from redthread.store import LocalStore, gitio


def _bare_remote(tmp_path):
    remote = tmp_path / "remote.git"
    subprocess.run(["git", "init", "-q", "--bare", "-b", "main", str(remote)], check=True)
    return remote


def _push_fresh_store(store: LocalStore, remote_url: str) -> None:
    root = store.layout.root
    gitio.configure_identity(root, "Test", "test@example.com")
    gitio.set_remote(root, remote_url)
    gitio.sync(root, "init store")


def _clone(remote, tmp_path, name: str) -> LocalStore:
    gitio.clone(str(remote), tmp_path / name)
    store = LocalStore(tmp_path / name)
    gitio.configure_identity(store.layout.root, "Test", "test@example.com")
    gitio.set_remote(store.layout.root, str(remote))
    return store


def test_ml_pipeline_present_matches_upstream_handoffs(tmp_path):
    remote = _bare_remote(tmp_path)
    store_a = LocalStore.init(
        tmp_path / "clone-a", project_id="ml-demo", phases=["train", "eval", "present"]
    )
    _push_fresh_store(store_a, str(remote))
    run = store_a.start_run()

    run_training(store_a, run.run_id)
    gitio.sync(store_a.layout.root, "train done")

    store_b = _clone(remote, tmp_path, "clone-b")
    run_eval(store_b, run.run_id)
    gitio.sync(store_b.layout.root, "eval done")

    store_c = _clone(remote, tmp_path, "clone-c")  # a THIRD node does present
    output_dir = tmp_path / "out"
    present_handoff = run_present(store_c, run.run_id, output_dir)
    gitio.sync(store_c.layout.root, "present done")

    train_handoff = store_c.get_handoff(run.run_id, "train")
    eval_handoff = store_c.get_handoff(run.run_id, "eval")
    assert present_handoff.key_results["phases_summarized"] == 2

    report = (output_dir / "report.md").read_text(encoding="utf-8")
    assert train_handoff.headline in report
    assert eval_handoff.headline in report
    assert f"| val_acc | {train_handoff.key_results['val_acc']} |" in report

    deck = Presentation(str(output_dir / "deck.pptx"))
    assert len(deck.slides) == 3  # title + train + eval
    slide_titles = [s.shapes.title.text for s in deck.slides]
    assert any(train_handoff.headline in t for t in slide_titles)
    assert any(eval_handoff.headline in t for t in slide_titles)

    index_md = (output_dir / "docs" / "index.md").read_text(encoding="utf-8")
    assert "train" in index_md
    assert "eval" in index_md

    record = store_c.get_run(run.run_id)
    assert record.phases["present"] == "done"


def test_app_pipeline_present_matches_upstream_handoffs(tmp_path):
    remote = _bare_remote(tmp_path)
    store_a = LocalStore.init(
        tmp_path / "clone-a", project_id="app-demo", phases=["build", "test", "present"]
    )
    _push_fresh_store(store_a, str(remote))
    run = store_a.start_run()

    artifact_source = tmp_path / "app.bin"
    artifact_source.write_bytes(b"compiled app")
    run_build(store_a, run.run_id, artifact_source)
    gitio.sync(store_a.layout.root, "build done")

    store_b = _clone(remote, tmp_path, "clone-b")
    run_tests(store_b, run.run_id)
    gitio.sync(store_b.layout.root, "test done")

    store_c = _clone(remote, tmp_path, "clone-c")
    output_dir = tmp_path / "out"
    present_handoff = run_present(store_c, run.run_id, output_dir)

    build_handoff = store_c.get_handoff(run.run_id, "build")
    test_handoff = store_c.get_handoff(run.run_id, "test")
    assert present_handoff.key_results["phases_summarized"] == 2

    report = (output_dir / "report.md").read_text(encoding="utf-8")
    assert build_handoff.headline in report
    assert test_handoff.headline in report
    assert "| coverage_pct | 87 |" in report

    deck = Presentation(str(output_dir / "deck.pptx"))
    assert len(deck.slides) == 3

    docs_build = (output_dir / "docs" / "build.md").read_text(encoding="utf-8")
    assert "enabled tree-shaking" in docs_build
